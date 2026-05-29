"""Generate PPTX report from JSON.

Usage: python gen_report_pptx.py <report.json> <output.pptx>
Reads report JSON from file, writes PPTX to output path.
"""

import json
import sys
import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Color palette — professional blue theme
PRIMARY = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xE5, 0x3E, 0x3E)
GREEN = RGBColor(0x38, 0xA1, 0x69)
YELLOW = RGBColor(0xD6, 0x9E, 0x2E)


def load_report(path):
    with open(path, 'r', encoding='utf-8') as f:
        return json.load(f)


def strip_tags(html):
    """Simple HTML tag stripper."""
    import re
    text = re.sub(r'<[^>]*>', '', html)
    text = text.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>').replace('&nbsp;', ' ')
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def add_slide_number(slide, num, total):
    """Add slide number footer."""
    left = Inches(0.5)
    top = Inches(7.1)
    width = Inches(9)
    height = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


def create_pptx(report, output_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    sections = report.get('sections', [])
    refs = report.get('references', [])
    meta = report.get('meta', {})

    # Pre-compute section chunks for accurate slide count
    chars_per_slide = 500
    section_chunks = []
    for section in sections:
        text = strip_tags(section.get('content', ''))
        chunks = []
        if len(text) <= chars_per_slide:
            chunks = [text]
        else:
            paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
            current = ""
            for para in paragraphs:
                if len(current) + len(para) < chars_per_slide:
                    current += para + "\n\n"
                else:
                    if current:
                        chunks.append(current.strip())
                    if len(para) > chars_per_slide:
                        words = para.split()
                        sub = ""
                        for w in words:
                            if len(sub) + len(w) < chars_per_slide:
                                sub += w + " "
                            else:
                                chunks.append(sub.strip())
                                sub = w + " "
                        if sub.strip():
                            chunks.append(sub.strip())
                    else:
                        current = para + "\n\n"
            if current.strip():
                chunks.append(current.strip())
        if not chunks:
            chunks = [text]
        section_chunks.append(chunks)

    total_content_slides = sum(len(chunks) for chunks in section_chunks)
    total_slides = 2 + total_content_slides + 1  # cover + toc + content + refs
    slide_num = 0

    # ── Cover Slide ──
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Blue accent bar
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0), Inches(0.12), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = report.get('title', '研究报告')
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8.5), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"{report.get('type', '深度报告')}  |  {meta.get('generatedAt', '')[:10]}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = GRAY

    # Stats summary
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8.5), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    parts = []
    if meta.get('sourceCount'):
        parts.append(f"{meta['sourceCount']} 条来源")
    if meta.get('citedCount'):
        parts.append(f"{meta['citedCount']} 条引用")
    if meta.get('verifiedClaims'):
        parts.append(f"{meta['verifiedClaims']} 条已验证主张")
    p3.text = "  ·  ".join(parts)
    p3.font.size = Pt(12)
    p3.font.color.rgb = GRAY

    add_slide_number(slide, slide_num, total_slides)

    # ── Table of Contents ──
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "目录"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Blue line
    shape = slide.shapes.add_shape(1, Inches(1), Inches(1.5), Inches(2), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, section in enumerate(sections):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = f"{i + 1}. {section['title']}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(12)

    add_slide_number(slide, slide_num, total_slides)

    # ── Content Slides ──
    for i, section in enumerate(sections):
        chunks = section_chunks[i]
        cit_count = len(section.get('citations', []))

        for chunk_idx, chunk in enumerate(chunks):
            slide_num += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

            # Section header with part number if multi-slide
            header_text = section['title']
            if len(chunks) > 1:
                header_text += f" ({chunk_idx + 1}/{len(chunks)})"

            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.5), Inches(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = header_text
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = PRIMARY

            # Separator line
            shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.2), Inches(8.5), Inches(0.02))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            shape.line.fill.background()

            # Content text
            txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.2))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = chunk
            p2.font.size = Pt(13)
            p2.font.color.rgb = DARK
            p2.line_spacing = Pt(20)

            # Citation count badge (only on last chunk slide)
            if chunk_idx == len(chunks) - 1 and cit_count > 0:
                txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(3), Inches(0.3))
                tf3 = txBox3.text_frame
                p3 = tf3.paragraphs[0]
                p3.text = f"引用来源: {cit_count} 条"
                p3.font.size = Pt(10)
                p3.font.color.rgb = GRAY

            add_slide_number(slide, slide_num, total_slides)

    # ── References Slide ──
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.5), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "参考文献"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK

    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, ref in enumerate(refs):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        cred = f"[{ref.get('credibility', 'B')}级]"
        p.text = f"[{ref['index']}] {ref['title']}  {cred}  {ref.get('url', '')}"
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p.space_after = Pt(6)

    add_slide_number(slide, slide_num, total_slides)

    prs.save(output_path)
    return output_path


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python gen_report_pptx.py <report.json> <output.pptx>")
        sys.exit(1)

    report_path = sys.argv[1]
    output_path = sys.argv[2]

    report = load_report(report_path)
    result = create_pptx(report, output_path)
    print(f"PPTX saved: {result}")
